import PyPDF2
from docx import Document
from pptx import Presentation
import io


def extract_text_from_file(file_obj):
    try:
        filename = file_obj.name.lower()
        
        if filename.endswith('.pdf'):
            return extract_text_from_pdf(file_obj)
        elif filename.endswith('.docx'):
            return extract_text_from_docx(file_obj)
        elif filename.endswith('.pptx'):
            return extract_text_from_pptx(file_obj)
        elif filename.endswith('.txt'):
            return extract_text_from_txt(file_obj)
        else:
            return f"Tipo de arquivo não suportado: {filename}"
            
    except Exception as e:
        return f"Erro ao extrair texto do arquivo: {str(e)}"


def extract_text_from_pdf(file_obj):
    try:
        pdf_reader = PyPDF2.PdfReader(file_obj)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip() if text.strip() else "Nenhum texto encontrado no PDF"
    except Exception as e:
        return f"Erro ao ler PDF: {str(e)}"


def extract_text_from_docx(file_obj):
    try:
        doc = Document(file_obj)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        return text.strip() if text.strip() else "Nenhum texto encontrado no DOCX"
    except Exception as e:
        return f"Erro ao ler DOCX: {str(e)}"


def extract_text_from_pptx(file_obj):
    try:
        prs = Presentation(file_obj)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip() if text.strip() else "Nenhum texto encontrado no PPTX"
    except Exception as e:
        return f"Erro ao ler PPTX: {str(e)}"


def extract_text_from_txt(file_obj):
    try:
        content = file_obj.read()
        if isinstance(content, bytes):
            text = content.decode('utf-8', errors='ignore')
        else:
            text = content
        return text.strip() if text.strip() else "Arquivo TXT vazio"
    except Exception as e:
        return f"Erro ao ler TXT: {str(e)}"
